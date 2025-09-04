#!/usr/bin/env python3
"""
SharePoint Weekly - Duplicate folder + rename (APP-ONLY, confidential client)

Qué hace:
1) Duplica una carpeta origen de SharePoint en el mismo padre,
   nombrando la nueva "W.E. %b %d %Y" (ej. "W.E. Sep 03 2025") en Europe/Madrid.
2) Renombra todos los archivos de primer nivel dentro de la carpeta nueva:
   - De: "11 Apr - EQT Project Updates.pptx"
   - A : "03 Sep - EQT Project Updates.pptx" (usando la fecha de hoy).

Autenticación:
- Microsoft Graph con **Application permissions** (app-only) usando client secret.
- Necesitas dar a la app Application Permissions en Graph:
  - Opción rápida: Sites.ReadWrite.All + Admin consent.
  - Opción segura: Sites.Selected + Admin consent + dar acceso al sitio vía Graph o PnP.

Config (.env):
TENANT_ID=xxxxxxxx-xxxx-xxxx-xxxx-xxxxxxxxxxxx
CLIENT_ID=xxxxxxxx-xxxx-xxxx-xxxx-xxxxxxxxxxxx
CLIENT_SECRET=your_super_secret
SP_HOSTNAME=techtorch.sharepoint.com
SP_SITE_PATH=/sites/DMTolls
SP_LIBRARY_NAME=Documents
SP_SOURCE_FOLDER_PATH=Delivery Excellence/MVP DeliveryHub/Delivery/TESTWeekly PE Updates/W.E. Sept 5 2025

Dependencias:
pip install msal python-dotenv requests pytz
"""

import os
import re
import time
import json
import sys
from datetime import datetime
from typing import Optional
from pptx import Presentation
from io import BytesIO
from pptx.util import Pt


import pytz
import requests
from msal import ConfidentialClientApplication
from dotenv import load_dotenv
from typing import Optional, Any
from datetime import datetime, timedelta



GRAPH = "https://graph.microsoft.com/v1.0"

# ------------------------ Utilidades de logging y ENV -------------------------

def log(msg: str) -> None:
    print(f"[{datetime.now().isoformat(timespec='seconds')}] {msg}", flush=True)

def require_env(name: str) -> str:
    v = os.getenv(name)
    if not v:
        raise RuntimeError(f"Missing required environment variable: {name}")
    return v

# ------------------------------- Autenticación -------------------------------

def get_token_app_only(tenant_id: str, client_id: str, client_secret: str) -> str:
    app = ConfidentialClientApplication(
        client_id=client_id,
        client_credential=client_secret,
        authority=f"https://login.microsoftonline.com/{tenant_id}",
    )
    result = app.acquire_token_for_client(scopes=["https://graph.microsoft.com/.default"])
    if "access_token" not in result:
        raise RuntimeError(f"Failed to acquire app-only token: {result}")
    return result["access_token"]

def gget(token: str, url: str, **kwargs) -> requests.Response:
    r = requests.get(url, headers={"Authorization": f"Bearer {token}"}, **kwargs)
    if not r.ok:
        raise RuntimeError(f"GET {url} -> {r.status_code} {r.text}")
    return r

def gpost(token: str, url: str, **kwargs) -> requests.Response:
    r = requests.post(url, headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"}, **kwargs)
    if r.status_code not in (200, 201, 202):
        raise RuntimeError(f"POST {url} -> {r.status_code} {r.text}")
    return r

def gpatch(token: str, url: str, **kwargs) -> requests.Response:
    r = requests.patch(url, headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"}, **kwargs)
    if not r.ok:
        raise RuntimeError(f"PATCH {url} -> {r.status_code} {r.text}")
    return r

def _next_friday(dt: datetime) -> datetime:
    # Monday=0 … Friday=4
    days_ahead = (4 - dt.weekday()) % 7
    if days_ahead == 0:
        days_ahead = 7  # if today is Friday, pick the Friday of next week
    return dt + timedelta(days=days_ahead)

# ------------------------ Resolución de sitio y drive ------------------------

def resolve_site_id(token: str, hostname: str, site_path: str) -> str:
    url = f"{GRAPH}/sites/{hostname}:{site_path}"
    r = gget(token, url)
    return r.json()["id"]

def find_drive_id(token: str, site_id: str, library_name: str) -> str:
    url = f"{GRAPH}/sites/{site_id}/drives"
    drives = gget(token, url).json().get("value", [])
    for d in drives:
        if d.get("name") == library_name:
            return d["id"]
    if library_name.strip().lower() == "shared documents":
        for d in drives:
            if d.get("name", "").lower() == "documents":
                return d["id"]
    for d in drives:
        if d.get("driveType") == "documentLibrary":
            return d["id"]
    raise RuntimeError(f"Drive '{library_name}' not found. Found: {[d.get('name') for d in drives]}")

def get_item_by_path(token: str, drive_id: str, path: str) -> dict:
    url = f"{GRAPH}/drives/{drive_id}/root:/{path}"
    return gget(token, url).json()

# ------------------------------ Copia de carpeta -----------------------------

def copy_folder_with_name(token: str, drive_id: str, source_item_id: str, parent_id: str, new_name: str) -> None:
    url = f"{GRAPH}/drives/{drive_id}/items/{source_item_id}/copy"
    payload = {
        "parentReference": {"driveId": drive_id, "id": parent_id},
        "name": new_name,
    }
    r = requests.post(
        url,
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
            "Prefer": "respond-async",
        },
        data=json.dumps(payload),
    )
    if r.status_code not in (200, 201, 202):
        raise RuntimeError(f"POST {url} -> {r.status_code} {r.text}")

def find_child_by_name(token: str, drive_id: str, parent_id: str, name: str) -> Optional[dict]:
    url = f"{GRAPH}/drives/{drive_id}/items/{parent_id}/children?$select=id,name,folder,file"
    children = gget(token, url).json().get("value", [])
    for c in children:
        if c.get("name") == name:
            return c
    return None

def poll_until_child_exists(token: str, drive_id: str, parent_id: str, name: str, timeout_s: int = 1200, poll_s: int = 4) -> dict:
    log(f"Waiting for copied folder '{name}' to appear…")
    t0 = time.time()
    while time.time() - t0 < timeout_s:
        child = find_child_by_name(token, drive_id, parent_id, name)
        if child:
            return child
        time.sleep(poll_s)
    raise TimeoutError(f"Folder '{name}' did not appear after {timeout_s} s.")

def list_children(token: str, drive_id: str, folder_id: str) -> list[dict]:
    url = f"{GRAPH}/drives/{drive_id}/items/{folder_id}/children?$select=id,name,folder,file"
    items = []
    while True:
        rj = gget(token, url).json()
        items.extend(rj.get("value", []))
        url = rj.get("@odata.nextLink")
        if not url:
            break
    return items

# ------------------------- Nombres y renombrado de files ---------------------

MONTH_ABBR = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]

def today_strings() -> tuple[str, str]:
    tz = pytz.timezone("Europe/Madrid")
    now = datetime.now(tz)
    target = _next_friday(now)  # <-- was: now

    day2 = f"{target.day:02d}"               # e.g., "05"
    mon_abbr = MONTH_ABBR[target.month - 1]  # e.g., "Sep"
    folder = f"W.E. {mon_abbr} {day2} {target.year}"  # "W.E. Sep 05 2025"
    file_prefix = f"{day2} {mon_abbr}"                # "05 Sep"
    return folder, file_prefix


# Regex que detecta prefijo con día + mes (acepta "Sep", "Sept", "September", etc.)
PREFIX_RE = re.compile(r"""(?ix)
^ \s*
(\d{1,2})                    # día
\s+
([A-Za-z]{3,9}\.?)           # mes abreviado/completo
\s*
[-–—]                        # guion o dash
\s+
(.+?)                        # resto
\s*$
""")

def build_new_filename(original_name: str, new_prefix: str) -> str:
    base, ext = os.path.splitext(original_name)
    m = PREFIX_RE.match(base)
    if m:
        rest = m.group(3)
        return f"{new_prefix} - {rest}{ext}"
    else:
        return f"{new_prefix} - {base}{ext}"

def rename_item_with_collision_retry(token: str, drive_id: str, item_id: str, new_name: str, max_tries: int = 5) -> str:
    base, ext = os.path.splitext(new_name)
    for n in range(max_tries):
        candidate = f"{base}{'' if n==0 else f' ({n})'}{ext}"
        url = f"{GRAPH}/drives/{drive_id}/items/{item_id}"
        r = requests.patch(url, headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
                           data=json.dumps({"name": candidate}))
        if r.ok:
            return r.json().get("name", candidate)
        if r.status_code != 409:
            raise RuntimeError(f"PATCH rename -> {r.status_code} {r.text}")
    raise RuntimeError(f"Could not rename after {max_tries} attempts.")

# ---------- Date formatting for slides ----------
def today_slide_date(prefer_sept_with_t: bool = True, tz_name: str = "Europe/Madrid") -> str:
    tz = pytz.timezone(tz_name)
    now = datetime.now(tz)
    target = _next_friday(now)  # <-- was: now

    day = f"{target.day:02d}"  # leading zero, e.g., "01"
    mon = MONTH_ABBR[target.month - 1]  # "Sep"
    if prefer_sept_with_t and mon == "Sep":
        mon = "Sept"
    return f"{day} {mon} {target.year}"  # "5 Sept 2025"


WEEK_ENDING_PATTERNS = [
    re.compile(r"(?i)\b(W\.?\s*E\.?\s*[:\-]?\s*)(\d{1,2}\s+[A-Za-z]{3,9}\.?\s+\d{4})\b"),
    re.compile(r"(?i)\b(Week\s*Ending\s*[:\-]?\s*)(\d{1,2}\s+[A-Za-z]{3,9}\.?\s+\d{4})\b"),
]

def replace_week_ending_text(text: str, new_date_label: str) -> tuple[str, bool]:
    changed = False
    new_text = text
    for rx in WEEK_ENDING_PATTERNS:
        new_text2, n = rx.subn(rf"\g<1>{new_date_label}", new_text)
        if n:
            changed = True
            new_text = new_text2
    return new_text, changed

    """
    Replace 'W.E. 5 Sept 2025' (or 'Week Ending 5 Sept 2025') with today's label.
    Returns (new_text, changed)
    """
    changed = False
    new_text = text
    for rx in WEEK_ENDING_PATTERNS:
        new_text, n = rx.subn(lambda m: f"{m.group(1)}{new_date_label}", new_text)
        if n:
            changed = True
    return new_text, changed

# ---------- Graph file download/upload ----------
def download_item_bytes(token: str, drive_id: str, item_id: str) -> bytes:
    url = f"{GRAPH}/drives/{drive_id}/items/{item_id}/content"
    r = requests.get(url, headers={"Authorization": f"Bearer {token}"}, allow_redirects=True)
    if not r.ok:
        raise RuntimeError(f"Download failed {r.status_code}: {r.text}")
    return r.content

def upload_item_bytes(token: str, drive_id: str, item_id: str, data: bytes) -> None:
    # Use upload session for anything non-trivial
    url = f"{GRAPH}/drives/{drive_id}/items/{item_id}/createUploadSession"
    sess = gpost(token, url, data=json.dumps({})).json()
    upload_url = sess["uploadUrl"]
    chunk = 5 * 1024 * 1024
    size = len(data)
    i = 0
    while i < size:
        j = min(i + chunk, size)
        headers = {
            "Content-Length": str(j - i),
            "Content-Range": f"bytes {i}-{j-1}/{size}",
        }
        rr = requests.put(upload_url, headers=headers, data=data[i:j])
        if rr.status_code not in (200, 201, 202):
            raise RuntimeError(f"Upload chunk failed {rr.status_code}: {rr.text}")
        i = j

# ---------- PPTX edit (slide 1 only) ----------
def update_pptx_first_slide_date(ppt_bytes: bytes, new_date_label: str) -> tuple[bytes, bool]:
    prs = Presentation(BytesIO(ppt_bytes))
    if len(prs.slides) == 0:
        return ppt_bytes, False

    s = prs.slides[0]
    changed_any = False

    for shape in s.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        tf = shape.text_frame

        # Full text across paragraphs (keeps \n so regex sees one block)
        original = "\n".join(p.text for p in tf.paragraphs)
        replaced, changed = replace_week_ending_text(original, new_date_label)

        if changed:
            # Snapshot current style before rewriting
            font_name, font_size, par_align, par_level = _snapshot_text_style(tf)
            _rewrite_textframe_preserving_style(tf, replaced, font_name, font_size, par_align, par_level)
            changed_any = True

    if not changed_any:
        return ppt_bytes, False

    out = BytesIO()
    prs.save(out)
    return out.getvalue(), True

def update_pptx_dates_in_folder(token: str, drive_id: str, folder_id: str) -> int:
    """
    Finds top-level *.pptx items in folder and updates slide 1 date.
    """
    target_label = today_slide_date(prefer_sept_with_t=True)  # e.g., "5 Sept 2025"
    log(f"Updating slide 1 date to: W.E. {target_label}")

    items = list_children(token, drive_id, folder_id)
    updated = 0
    for it in items:
        name = it.get("name", "")
        if "file" in it and name.lower().endswith(".pptx"):
            try:
                b = download_item_bytes(token, drive_id, it["id"])
                nb, changed = update_pptx_first_slide_date(b, target_label)
                if changed:
                    upload_item_bytes(token, drive_id, it["id"], nb)
                    log(f"Updated slide 1 date in: {name}")
                    updated += 1
                else:
                    log(f"No date string found on slide 1 (unchanged): {name}")
            except Exception as e:
                log(f"ERROR updating {name}: {e}")
    return updated

def _snapshot_text_style(tf) -> tuple[Optional[str], Optional[Any], Optional[Any], Optional[int]]:
    """
    Returns (font_name, font_size, first_par_align, first_par_level) from the first non-empty run/paragraph.
    Falls back to None if not explicitly set (theme-based). We'll default later to Arial / 24 pt if missing.
    """
    font_name = None
    font_size = None
    par_align = None
    par_level = None

    for p in tf.paragraphs:
        if par_align is None:
            par_align = p.alignment
            try:
                par_level = p.level
            except Exception:
                par_level = None
        for r in p.runs:
            if r.text:
                if r.font.name:
                    font_name = r.font.name
                if r.font.size:
                    font_size = r.font.size
                if font_name and font_size:
                    break
        if font_name and font_size:
            break
    return font_name, font_size, par_align, par_level

def _rewrite_textframe_preserving_style(tf, text: str, font_name, font_size, par_align, par_level):
    """
    Clears the text frame and writes `text` (may contain \n) while restoring font family/size
    and basic paragraph props. Defaults to Arial / 24 pt if originals are not explicitly set.
    """
    default_name = font_name or "Arial"
    default_size = font_size or Pt(24)

    tf.clear()
    lines = text.split("\n")

    def _format_paragraph(p):
        if par_align is not None:
            p.alignment = par_align
        try:
            if par_level is not None:
                p.level = par_level
        except Exception:
            pass
        # Ensure at least one run exists (setting p.text creates one)
        r = p.runs[0]
        r.font.name = default_name
        r.font.size = default_size

    # first line
    p0 = tf.paragraphs[0]
    p0.text = lines[0]
    _format_paragraph(p0)

    # remaining lines
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        _format_paragraph(p)

# ------------------------------ Mail helpers ---------------------------------

def get_item_fields(token: str, drive_id: str, item_id: str, select: str = "id,name,webUrl") -> dict:
    url = f"{GRAPH}/drives/{drive_id}/items/{item_id}?$select={select}"
    return gget(token, url).json()

def _addr_list(csv_or_list) -> list[dict]:
    if not csv_or_list:
        return []
    if isinstance(csv_or_list, str):
        addrs = [a.strip() for a in csv_or_list.split(",") if a.strip()]
    else:
        addrs = [str(a).strip() for a in csv_or_list if str(a).strip()]
    return [{"emailAddress": {"address": a}} for a in addrs]

def send_mail_app_only(
    token: str,
    sender_user_upn: str,
    to: list[str] | str,
    subject: str,
    html_body: str,
    cc: list[str] | str | None = None,
    bcc: list[str] | str | None = None,
    save_to_sent: bool = True,
) -> None:
    """
    Sends an HTML email as `sender_user_upn` using Microsoft Graph (app-only).
    Requires Mail.Send (Application) with admin consent; sender mailbox must exist.
    """
    url = f"{GRAPH}/users/{sender_user_upn}/sendMail"
    payload = {
        "message": {
            "subject": subject,
            "body": {"contentType": "HTML", "content": html_body},
            "toRecipients": _addr_list(to),
        },
        "saveToSentItems": bool(save_to_sent),
    }
    cc_list = _addr_list(cc)
    if cc_list:
        payload["message"]["ccRecipients"] = cc_list
    bcc_list = _addr_list(bcc)
    if bcc_list:
        payload["message"]["bccRecipients"] = bcc_list

    gpost(token, url, data=json.dumps(payload))


# --------------------------------- Main --------------------------------------

def main():
    load_dotenv()

    tenant_id = require_env("TENANT_ID")
    client_id = require_env("CLIENT_ID")
    client_secret = require_env("CLIENT_SECRET")
    sp_hostname = require_env("SP_HOSTNAME")
    sp_site_path = require_env("SP_SITE_PATH")
    sp_library = require_env("SP_LIBRARY_NAME")
    source_folder_path = require_env("SP_SOURCE_FOLDER_PATH")

    new_folder_name, file_prefix = today_strings()
    token = get_token_app_only(tenant_id, client_id, client_secret)

    log("Resolving site…")
    site_id = resolve_site_id(token, sp_hostname, sp_site_path)
    log(f"Site id: {site_id}")

    log("Finding library…")
    drive_id = find_drive_id(token, site_id, sp_library)
    log(f"Drive id: {drive_id}")

    log("Locating source folder…")
    source_item = get_item_by_path(token, drive_id, source_folder_path)
    if source_item.get("folder") is None:
        raise RuntimeError("SP_SOURCE_FOLDER_PATH is not a folder.")
    parent_id = source_item["parentReference"]["id"]
    source_id = source_item["id"]
    log(f"Source folder id: {source_id} (parent {parent_id})")

    log(f"Copying folder as '{new_folder_name}'…")
    copy_folder_with_name(token, drive_id, source_id, parent_id, new_folder_name)

    new_folder_item = poll_until_child_exists(token, drive_id, parent_id, new_folder_name)
    new_folder_id = new_folder_item["id"]
    log(f"New folder ready: {new_folder_name} (id {new_folder_id})")

    log("Renaming files in new folder (top-level only)…")
    children = list_children(token, drive_id, new_folder_id)
    renamed = 0
    for it in children:
        if "file" in it:
            original = it["name"]
            new_name = build_new_filename(original, file_prefix)
            if new_name != original:
                applied = rename_item_with_collision_retry(token, drive_id, it["id"], new_name)
                log(f"Renamed file: {original} -> {applied}")
                renamed += 1
            else:
                log(f"No change needed: {original}")
        else:
            log(f"Skipping folder: {it.get('name')}")
    log(f"Done. Files renamed: {renamed}")
        # ... after: log(f"Done. Files renamed: {renamed}")
    log("Updating slide 1 date in PPTX files…")
    updated = update_pptx_dates_in_folder(token, drive_id, new_folder_id)
    log(f"Done. PPTX files updated: {updated}")

    # -------------------------- Send notification email --------------------------
    sender_upn = require_env("MAIL_SENDER_UPN")
    to_csv    = require_env("MAIL_TO")
    cc_csv    = os.getenv("MAIL_CC", "")
    bcc_csv   = os.getenv("MAIL_BCC", "")

    # Get a friendly link to the newly created W.E. folder
    nf_meta = get_item_fields(token, drive_id, new_folder_id, select="id,name,webUrl")
    folder_name = nf_meta.get("name", new_folder_name)
    folder_link = nf_meta.get("webUrl", "(link unavailable)")

    subject = f"PE Weekly Updates {folder_name}"
    html_body = f"""
    <html>
      <body style="font-family: Segoe UI, Arial, sans-serif;">
        <p>Hi team,</p>
        <p>Please see linked below the folder where all the ✨new PPTs ✨ for this week live! Please find the slides for your clients, organized by PE firm, and make all updates for progress from this week. Let me know when updates are completed – I will be aiming to send PDFs to Jordi on Friday afternoon EST.</p>
        <p><a href="{folder_link}">{folder_name}</a></p>
        <ol>
          <li><span style="background-color: #f0f0f0;">@Samantha Starr</span> can you please update <b>Benevity, Avetta, Calabrio & TurboTax/Bain</b></li>
          <li><span style="background-color: #f0f0f0;">@Ash Shah</span> please confirm updates (if any) for <b>HVD</b> and <b>ClearCourse</b></li>
          <li><span style="background-color: #f0f0f0;">@Tim Hegwood</span> can you please provide updates for <b>Intel471</b></li>
          <li><span style="background-color: #f0f0f0;">@Marina Rodriguez</span> can you please update for <b>Waystone Data</b></li>
          <li><span style="background-color: #f0f0f0;">@Marta Ulanecka, @Agustin Russo, @Mikolaj Maslanka</span> can you please update <b>Logex</b></li>
          <li><span style="background-color: #f0f0f0;">@Frank Brugnot</span> please provide updates for <b>EQS</b></li>
          <li><span style="background-color: #f0f0f0;">@Julio Saiz</span> please provide updates for <b>USU</b></li>
          <li><span style="background-color: #f0f0f0;">@Luis Guillermo Salazar Madriz</span> please provide update for <b>Hypergene</b></li>
          <li><span style="background-color: #f0f0f0;">@Marta Ulanecka</span> can you please provide updates for <b>Lunanet, Mercell, AEA</b></li>
          <li><span style="background-color: #f0f0f0;">@Jake Harris</span> please provide updates for <b>Caseware</b></li>
          <li><span style="background-color: #f0f0f0;">@Nathan Lewis</span> please provide updates for <b>Suvoda</b></li>
          <li><span style="background-color: #f0f0f0;">@Kurt Mansperger</span> can you please update for <b>Magnet Forensics</b></li>
          <li><span style="background-color: #f0f0f0;">@Liz Shelton</span> can you please provide updates for <b>BillTrust</b></li>
        </ol>
        <p><span style="font-weight: bold; text-decoration: underline;">Reminders:</span></p>
        <ol>
          <li><b>Please review + follow the guidelines linked <a href="https://techtorch.sharepoint.com/:p:/r/sites/DMTolls/_layouts/15/Doc.aspx?sourcedoc=%7B07A9793F-4E37-456B-AB9E-F516F4D49AC3%7D&file=Update%20Guidelines.pptx&action=edit&mobileredirect=true">HERE</a> for consistency!</b></li>
          <li>Reminder to <span style="background-color: yellow;">please complete your updates on the slides & send to client - Due Thursday EOD</span></li> 
          <ol>
            <li>Save <span style="text-decoration: underline;">each slide on its own as a PDF</span></li>
            <li>Send a <span style="text-decoration: underline;">short email to the main contact/project sponsor</span> at your client with the PDF of the single slide attached (CC Ash for EU/Miguel for US)</li>
            <li>REMINDER: The goal is <b>not</b> to seek validation from clients, but rather to <b>ensure they are informed and prepared to provide additional context and next steps</b> internally if needed.</li>
          </ol>
          <li><span style="background-color: yellow;">PLEASE notify me via email/teams that your updates are done and you have sent the slide to the clients - <b>Due Friday 12pm EST</b></span></li>
        </ol>
        <p>Let me know if you have any questions!</p>
        <p>Thank you all!,<br/>Erica Charles</p>
      </body>
    </html>
    """

    log(f"Sending notification email to: {to_csv} (cc: {cc_csv or '—'}, bcc: {bcc_csv or '—'})")
    send_mail_app_only(
        token=token,
        sender_user_upn=sender_upn,
        to=to_csv,
        subject=subject,
        html_body=html_body,
        cc=cc_csv,
        bcc=bcc_csv,
        save_to_sent=True,
    )
    log("Email sent.")



if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        log(f"ERROR: {e}")
        sys.exit(1)
